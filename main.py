import os
import uuid
from pathlib import Path
from openai import OpenAI

from fastapi import FastAPI, UploadFile, File
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from pptx import Presentation

from fastapi.middleware.cors import CORSMiddleware

app = FastAPI(title="PPT Translator")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

client = OpenAI(
    api_key=os.getenv("ARK_API_KEY"),
    base_url="https://ark.cn-beijing.volces.com/api/v3",
)

BASE_DIR = Path(__file__).parent
UPLOAD_DIR = BASE_DIR / "uploads"
RESULT_DIR = BASE_DIR / "results"
STATIC_DIR = BASE_DIR / "static"

UPLOAD_DIR.mkdir(exist_ok=True)
RESULT_DIR.mkdir(exist_ok=True)
STATIC_DIR.mkdir(exist_ok=True)

@app.get("/api/ping")
def ping():
    return {"status": "ok", "message": "server is running"}
    
@app.get("/", response_class=HTMLResponse)
def home():
    index_path = STATIC_DIR / "index.html"

    if index_path.exists():
        return index_path.read_text(encoding="utf-8")

    return """
    <html>
        <body>
            <h2>PPT Translator is running.</h2>
            <p>Please create static/index.html.</p>
        </body>
    </html>
    """


@app.post("/api/translate-ppt")
async def translate_ppt(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pptx"):
        return JSONResponse(
            status_code=400,
            content={
                "status": "failed",
                "message": "Only .pptx files are supported."
            }
        )

    original_name = Path(file.filename).stem
    safe_name = original_name.replace(" ", "_")
    
    file_id = str(uuid.uuid4())
    
    input_filename = f"{file_id}.pptx"
    output_filename = f"{safe_name}_en.pptx"
    
    input_path = UPLOAD_DIR / input_filename
    output_path = RESULT_DIR / output_filename
    
    content = await file.read()
    input_path.write_bytes(content)

    try:
        prs = Presentation(str(input_path))

        text_items = []
        shape_map = {}
        
        MAX_TRANSLATE_COUNT = 80
        
        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    original_text = shape.text.strip()
        
                    if original_text:
                        item_id = f"slide_{slide_index}_shape_{shape_index}"
        
                        if len(text_items) < MAX_TRANSLATE_COUNT:
                            text_items.append({
                                "id": item_id,
                                "text": original_text
                            })
        
                            shape_map[item_id] = {
                                "shape": shape,
                                "source_text": original_text
                            }
        
        translations = batch_translate_by_doubao(text_items)
        
        for item_id, info in shape_map.items():
            shape = info["shape"]
            source_text = info["source_text"]
        
            translated_text = translations.get(item_id, source_text)
        
            shape.text = translated_text
            simple_reduce_font(shape, source_text, translated_text)

        prs.save(str(output_path))

    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={
                "status": "failed",
                "message": f"PPT processing failed: {str(e)}"
            }
        )

    return {
        "status": "success",
        "download_url": f"/api/download/{output_filename}"
    }


@app.get("/api/download/{filename}")
def download_file(filename: str):
    file_path = RESULT_DIR / filename

    if not file_path.exists():
        return JSONResponse(
            status_code=404,
            content={
                "status": "failed",
                "message": "File not found."
            }
        )

    return FileResponse(
        path=str(file_path),
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


def batch_translate_by_doubao(text_items):
    """
    批量调用豆包翻译 PPT 文本。
    text_items 格式：
    [
        {"id": "slide_0_shape_1", "text": "研究背景"},
        {"id": "slide_0_shape_2", "text": "本文提出..."}
    ]
    """

    if not text_items:
        return {}

    model = os.getenv("ARK_MODEL")

    if not os.getenv("ARK_API_KEY"):
        return {item["id"]: "[ARK_API_KEY_MISSING] " + item["text"] for item in text_items}

    if not model:
        return {item["id"]: "[ARK_MODEL_MISSING] " + item["text"] for item in text_items}

    input_json = {
        "items": text_items
    }

    prompt = f"""
你是一个专业的 PPT 中英翻译助手。
请将输入 JSON 中每个 item 的中文 text 翻译成英文。

要求：
1. 保持原意，不要扩写；
2. 译文适合放在 PPT 中，尽量简洁；
3. 保留数字、公式、变量名、英文缩写；
4. 如果原文已经是英文、数字或公式，可以保持不变；
5. 必须保留每个 item 的 id；
6. 只返回 JSON，不要返回解释文字；
7. 返回格式必须如下：

{{
  "items": [
    {{
      "id": "slide_0_shape_1",
      "translated_text": "Research Background"
    }}
  ]
}}

待翻译 JSON：
{json.dumps(input_json, ensure_ascii=False)}
"""

    try:
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "user", "content": prompt}
            ],
            temperature=0.2,
        )

        content = response.choices[0].message.content.strip()

        # 去掉模型可能返回的 ```json 包裹
        if content.startswith("```"):
            content = content.replace("```json", "").replace("```", "").strip()

        result = json.loads(content)

        translations = {}

        for item in result.get("items", []):
            item_id = item.get("id")
            translated_text = item.get("translated_text")

            if item_id and translated_text:
                translations[item_id] = translated_text

        return translations

    except Exception as e:
        print("Batch Doubao translation failed:", str(e))
        return {item["id"]: item["text"] for item in text_items}

def simple_reduce_font(shape, source_text: str, translated_text: str):
    """
    最简单的字号缩小规则。
    如果英文比中文长很多，就把字号稍微缩小。
    """
    ratio = len(translated_text) / max(len(source_text), 1)

    if ratio > 3:
        reduce_size = 6
    elif ratio > 2:
        reduce_size = 4
    elif ratio > 1.5:
        reduce_size = 2
    else:
        reduce_size = 0

    if reduce_size == 0:
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                new_size = run.font.size.pt - reduce_size
                if new_size < 10:
                    new_size = 10
                run.font.size = int(new_size * 12700)
